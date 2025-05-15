import logging
from pathlib import Path
import pytesseract
from PyPDF2 import PdfReader
from PIL import Image
from docx import Document
from pptx import Presentation
import pandas as pd
from project import settings
from pdfminer.high_level import extract_text as pdfminer_extract_text
import cv2
import spacy
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_qdrant import QdrantVectorStore
from qdrant_client import QdrantClient
from qdrant_client.http.models import Distance, VectorParams, MatchValue, Filter, FieldCondition
from langchain.schema import Document as LangChainDocument
import json
import xml.etree.ElementTree as ET
from langchain.prompts import PromptTemplate
from concurrent.futures import ThreadPoolExecutor
from typing import List, Optional, Dict
from .models import Document, OpenAIKey, DocumentAlert, Tenant, VectorStore, DocumentAccess
from knox.models import AuthToken
from django.shortcuts import get_object_or_404
import os
import subprocess
import platform
import speech_recognition as sr
from pydub import AudioSegment
import tempfile
from moviepy.editor import VideoFileClip
import whisper
import requests
from urllib.parse import urlparse
from functools import lru_cache
from collections import defaultdict
import nltk
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
import time
from langchain.chains.llm import LLMChain
from bs4 import BeautifulSoup
import yaml
import configparser

# Initialize logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Download NLTK data
nltk.download('punkt', quiet=True)
nltk.download('stopwords', quiet=True)

nlp = spacy.load("en_core_web_sm")

# Cache Qdrant client
@lru_cache(maxsize=1)
def get_qdrant_client():
    return QdrantClient(host="localhost", port=6333)

@lru_cache(maxsize=128)
def get_openai_api_key(user_id: int) -> str:
    try:
        if user_id:
            openai_key = OpenAIKey.objects.filter(user_id=user_id, is_valid=True).first()
            return openai_key.api_key if openai_key else settings.OPENAI_API_KEY
        return settings.OPENAI_API_KEY
    except Exception as e:
        logger.error(f"Error retrieving OpenAI API key for user_id {user_id}: {e}")
        return settings.OPENAI_API_KEY

def get_authenticated_user(token: str):
    try:
        auth_token = get_object_or_404(AuthToken, token_key=token)
        user = auth_token.user
        if not hasattr(user, 'tenant') or user.tenant is None:
            logger.error(f"User {user.username} has no associated tenant")
            raise ValueError("User has no associated tenant")
        logger.info(f"Authenticated user: {user.username}")
        return user
    except Exception as e:
        logger.error(f"Error authenticating user with token {token}: {e}")
        raise

def initialize_qdrant_collection(collection_name: str) -> QdrantClient:
    client = get_qdrant_client()
    try:
        client.get_collection(collection_name=collection_name)
        logger.info(f"Collection '{collection_name}' already exists.")
    except Exception:
        logger.info(f"Creating collection '{collection_name}'...")
        client.create_collection(
            collection_name=collection_name,
            vectors_config=VectorParams(size=1536, distance=Distance.COSINE)
        )
        logger.info(f"Collection '{collection_name}' created successfully.")
    return client

def get_qdrant_vector_store(user, collection_name: str) -> QdrantVectorStore:
    try:
        api_key = get_openai_api_key(user.id if user else 0)
        embedding_model = OpenAIEmbeddings(api_key=api_key)
        client = initialize_qdrant_collection(collection_name)
        return QdrantVectorStore(
            client=client,
            collection_name=collection_name,
            embedding=embedding_model,
        )
    except Exception as e:
        logger.error(f"Error initializing Qdrant vector store for collection {collection_name}: {e}")
        raise

def process_file(file=None, s3_file_url=None, file_name=None) -> tuple:
    try:
        if file:
            file_name = file_name or file.name
            with tempfile.NamedTemporaryFile(delete=False, suffix=Path(file_name).suffix) as tmp:
                for chunk in file.chunks():
                    tmp.write(chunk)
                tmp_path = tmp.name
        elif s3_file_url:
            parsed_url = urlparse(s3_file_url)
            file_name = file_name or parsed_url.path.split("/")[-1]
            response = requests.get(s3_file_url)
            if response.status_code != 200:
                raise ValueError(f"Failed to download file from S3 URL: {response.status_code}")
            with tempfile.NamedTemporaryFile(delete=False, suffix=Path(file_name).suffix) as tmp:
                tmp.write(response.content)
                tmp_path = tmp.name
        else:
            raise ValueError("Either file or s3_file_url must be provided.")
        logger.info(f"Processed file: {file_name}, temp path: {tmp_path}")
        return tmp_path, file_name
    except Exception as e:
        logger.error(f"Error processing file: {e}")
        raise

def extract_text_from_file(file_path: str, original_file_name: str) -> str:
    ext = Path(original_file_name).suffix.lower()
    logger.info(f"Starting text extraction for {original_file_name} with extension {ext}")
    
    extraction_methods = {
        ".pdf": extract_text_from_pdf,
        ".docx": extract_text_from_docx,
        ".png": extract_text_from_image,
        ".jpg": extract_text_from_image,
        ".jpeg": extract_text_from_image,
        ".pptx": extract_text_from_pptx,
        ".xls": extract_structured_from_excel_or_csv,
        ".xlsx": extract_structured_from_excel_or_csv,
        ".csv": extract_structured_from_excel_or_csv,
        ".txt": extract_text_from_txt,
        ".json": extract_structured_from_json,
        ".xml": extract_text_from_xml,
        ".html": extract_text_from_html,
        ".md": extract_text_from_md,
        ".yaml": extract_text_from_yaml,
        ".yml": extract_text_from_yaml,
        ".ini": extract_text_from_ini,
        ".cfg": extract_text_from_ini,
        ".ppt": extract_text_from_ppt,
        ".mp4": extract_text_from_video,
        ".avi": extract_text_from_video,
        ".mov": extract_text_from_video,
        ".mp3": extract_text_from_audio,
        ".wav": extract_text_from_audio,
        ".m4a": extract_text_from_audio,
    }
    
    method = extraction_methods.get(ext)
    if not method:
        raise ValueError(f"Unsupported file type: {ext}")
    
    extracted_data = method(file_path)
    if isinstance(extracted_data, (dict, list)):
        extracted_data = json.dumps(extracted_data)
    
    logger.info(f"Successfully extracted data from {original_file_name}: {type(extracted_data)}")
    return extracted_data.strip()

def extract_text_from_pdf(file_path: str) -> str:
    try:
        reader = PdfReader(file_path)
        if reader.is_encrypted:
            raise ValueError("The uploaded PDF is password-protected. Please upload an unlocked PDF.")
        text = pdfminer_extract_text(file_path)
        if not text.strip():
            raise ValueError("The PDF appears to be empty or unreadable after extraction.")
        return text.strip()
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {e}")
        raise

def extract_text_from_docx(file_path: str) -> str:
    try:
        doc = Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                full_text.append('\t'.join(row_data))
        return '\n'.join(full_text).strip()
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {e}")
        raise

def extract_text_from_image(file_path: str) -> str:
    try:
        img = cv2.imread(file_path)
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        blur = cv2.GaussianBlur(gray, (1, 1), 0)
        thresh = cv2.threshold(blur, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)[1]
        text = pytesseract.image_to_string(thresh)
        return text.strip()
    except Exception as e:
        logger.error(f"Error extracting text from image: {e}")
        raise

def extract_text_from_pptx(file_path: str) -> str:
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                text += slide.notes_slide.notes_text_frame.text + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {e}")
        raise

def extract_text_from_ppt(file_path: str) -> str:
    pptx_path = file_path + "x"
    try:
        logger.info(f"Converting PPT to PPTX on {platform.system()}...")
        if platform.system() == "Windows":
            import comtypes.client
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1
            ppt = powerpoint.Presentations.Open(file_path, WithWindow=False)
            ppt.SaveAs(pptx_path, 24)
            ppt.Close()
            powerpoint.Quit()
        else:
            subprocess.run(["libreoffice", "--headless", "--convert-to", "pptx", file_path, "--outdir", os.path.dirname(file_path)], check=True)
        return extract_text_from_pptx(pptx_path)
    except Exception as e:
        logger.error(f"Error extracting text from PPT: {e}")
        return ""
    finally:
        if os.path.exists(pptx_path):
            os.remove(pptx_path)

def extract_structured_from_excel_or_csv(file_path: str):
    try:
        ext = Path(file_path).suffix.lower()
        if ext == ".csv":
            df = pd.read_csv(file_path, encoding="utf-8", errors='replace')
            return df.to_dict(orient="records")
        elif ext in [".xlsx", ".xls"]:
            engine = "openpyxl" if ext == ".xlsx" else "xlrd"
            xls = pd.ExcelFile(file_path, engine=engine)
            dfs = {}
            for sheet_name in xls.sheet_names:
                sheet_df = xls.parse(sheet_name)
                for col in sheet_df.columns:
                    if pd.api.types.is_datetime64_any_dtype(sheet_df[col]):
                        sheet_df[col] = sheet_df[col].astype(str)
                dfs[sheet_name] = sheet_df.to_dict(orient="records")
            return dfs
        else:
            raise ValueError(f"Unsupported Excel/CSV file extension: {ext}")
    except Exception as e:
        logger.error(f"Error extracting structured data from Excel/CSV: {e}")
        raise

def extract_text_from_txt(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read().strip()
    except Exception as e:
        logger.error(f"Error extracting text from TXT: {e}")
        raise

def extract_structured_from_json(file_path: str):
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return data
    except Exception as e:
        logger.error(f"Error extracting structured data from JSON: {e}")
        raise

def extract_text_from_xml(file_path: str) -> str:
    try:
        tree = ET.parse(file_path)
        root = tree.getroot()
        return ET.tostring(root, encoding='utf-8').decode('utf-8').strip()
    except Exception as e:
        logger.error(f"Error extracting text from XML: {e}")
        raise

def extract_text_from_html(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f, "html.parser")
        return soup.get_text().strip()
    except Exception as e:
        logger.error(f"Error extracting text from HTML: {e}")
        raise

def extract_text_from_md(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read().strip()
    except Exception as e:
        logger.error(f"Error extracting text from Markdown: {e}")
        raise

def extract_text_from_yaml(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            data = yaml.safe_load(f)
        return yaml.dump(data).strip()
    except Exception as e:
        logger.error(f"Error extracting text from YAML: {e}")
        raise

def extract_text_from_ini(file_path: str) -> str:
    try:
        config = configparser.ConfigParser()
        config.read(file_path)
        output = []
        for section in config.sections():
            output.append(f"[{section}]")
            for key, value in config.items(section):
                output.append(f"{key} = {value}")
        return "\n".join(output).strip()
    except Exception as e:
        logger.error(f"Error extracting text from INI: {e}")
        raise

def extract_text_from_video(file_path: str) -> str:
    try:
        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp_audio:
            video = VideoFileClip(file_path)
            audio = video.audio
            audio.write_audiofile(tmp_audio.name, logger=None)
            audio.close()
            video.close()
            return extract_text_from_audio(tmp_audio.name)
    except Exception as e:
        logger.error(f"Error extracting text from video: {e}")
        return ""
    finally:
        if 'tmp_audio' in locals() and os.path.exists(tmp_audio.name):
            os.remove(tmp_audio.name)

def extract_text_from_audio(file_path: str) -> str:
    try:
        model = whisper.load_model("base")
        result = model.transcribe(file_path, language="en")
        return result["text"].strip()
    except Exception as e:
        logger.error(f"Error extracting text from audio: {e}")
        return ""

def extract_keywords(query: str) -> List[str]:
    try:
        stop_words = set(stopwords.words('english'))
        tokens = word_tokenize(query.lower())
        keywords = [token for token in tokens if token.isalnum() and token not in stop_words]
        doc = nlp(query)
        entities = [ent.text.lower() for ent in doc.ents if ent.label_ in ["ORG", "GPE", "PERSON", "DATE"]]
        return list(set(keywords + entities))
    except Exception as e:
        logger.error(f"Error extracting keywords: {e}")
        return []

def hybrid_search(
    query: str,
    user,
    collection_name: str,
    vector_store: QdrantVectorStore,
    top_k: int = 5,
    vector_weight: float = 0.7,
    keyword_weight: float = 0.3
) -> List[Dict]:
    logger.info(f"Performing hybrid search for query: {query}")
    try:
        keywords = extract_keywords(query)
        logger.info(f"Extracted keywords: {keywords}")
        
        retriever = vector_store.as_retriever(search_kwargs={"k": top_k * 2})
        vector_docs = retriever.get_relevant_documents(query)
        
        client = get_qdrant_client()
        keyword_results = []
        if keywords:
            keyword_filter = Filter(
                should=[FieldCondition(key="page_content", match=MatchValue(value=kw)) for kw in keywords]
            )
            keyword_results = client.scroll(
                collection_name=collection_name,
                scroll_filter=keyword_filter,
                limit=top_k * 2,
                with_vectors=False,
                with_payload=True
            )[0] or []
        
        doc_scores = defaultdict(float)
        doc_map = {}
        
        for i, doc in enumerate(vector_docs):
            doc_id = doc.metadata.get("document_id", str(i))
            score = 1.0 / (i + 1)
            doc_scores[doc_id] += score * vector_weight
            doc_map[doc_id] = {"content": doc.page_content, "metadata": doc.metadata}
        
        for i, record in enumerate(keyword_results):
            doc_id = record.payload.get("metadata", {}).get("document_id", str(i))
            score = 1.0 / (i + 1)
            doc_scores[doc_id] += score * keyword_weight
            if doc_id not in doc_map:
                doc_map[doc_id] = {
                    "content": record.payload.get("page_content", ""),
                    "metadata": record.payload.get("metadata", {})
                }
        
        sorted_docs = sorted(doc_scores.items(), key=lambda x: x[1], reverse=True)[:top_k]
        results = [doc_map[doc_id] for doc_id, _ in sorted_docs]
        
        logger.info(f"Hybrid search returned {len(results)} documents")
        return results
    except Exception as e:
        logger.error(f"Error in hybrid search: {e}")
        raise

def ask_question(
    query: str,
    source_type: str,
    documents: Optional[List[dict]] = None,
    user=None,
    collection_name: Optional[str] = None,
    assistant_instructions: Optional[str] = None
) -> str:
    try:
        start = time.time()
        api_key = get_openai_api_key(user.id if user else 0)
        llm = ChatOpenAI(api_key=api_key, temperature=0)

        # Use assistant_instructions if provided; otherwise, use default system prompt
        default_system_prompt = """
        You are a helpful document analyst. Provide accurate, concise, and detailed answers based on the provided document context.
        Use the information in the documents to answer the question directly and avoid including irrelevant details.
        If the documents do not contain enough information, state so clearly.
        """
        system_prompt = assistant_instructions if assistant_instructions else default_system_prompt

        final_context = ""
        max_token_threshold = 4000
        batch_size = 5

        if documents:
            logger.info("Using provided documents as context")
            grouped_docs = defaultdict(list)
            for doc in documents:
                document_id = doc["metadata"].get("document_id", "unknown")
                grouped_docs[document_id].append(doc)

            mini_summaries = []
            for document_id, docs in grouped_docs.items():
                combined_text = "\n".join([d["content"] for d in docs])
                est_tokens = len(combined_text) // 4
                if est_tokens <= max_token_threshold:
                    mini_summary = summarize_context(combined_text, user)
                else:
                    batches = list(batch_documents(
                        [LangChainDocument(page_content=d["content"]) for d in docs],
                        batch_size=batch_size
                    ))
                    batch_summaries = []
                    with ThreadPoolExecutor(max_workers=5) as executor:
                        futures = [executor.submit(summarize_batch, batch) for batch in batches]
                        for future in futures:
                            batch_summaries.append(future.result())
                    combined_summary = "\n".join(batch_summaries)
                    mini_summary = summarize_context(combined_summary, user)
                mini_summaries.append(mini_summary)

            final_context = "\n\n".join(mini_summaries)
        else:
            logger.info("Using hybrid search for context")
            if not collection_name and user:
                collection_name = user.tenant.collection_name if hasattr(user, 'tenant') and user.tenant else None
            if not collection_name:
                raise ValueError("Collection name is required for hybrid search")
            vector_store = get_qdrant_vector_store(user, collection_name)
            retrieved_docs = hybrid_search(query, user, collection_name, vector_store)
            final_context = "\n".join([doc["content"] for doc in retrieved_docs])

        end = time.time()
        logger.info(f"Retrieved and processed context in: {end - start:.2f} seconds")

        prompt_template = PromptTemplate(
            template="""{system_prompt}\n\nContext:\n{context}\n\nQuestion: {question}\n\nAnswer:""",
            input_variables=["system_prompt", "context", "question"]
        )

        qa_chain = LLMChain(
            llm=llm,
            prompt=prompt_template.partial(system_prompt=system_prompt, context=final_context)
        )

        start_time = time.time()
        result = qa_chain.invoke({"question": query})["text"]
        end_time = time.time()
        logger.info(f"Question answered in: {end_time - start_time:.2f} seconds")
        return result
    except Exception as e:
        logger.error(f"Error in ask_question: {e}")
        raise

# def retrieve_documents_by_vector_id(vector_store_id: str, user, collection_name: str) -> List[Dict]:
#     try:
#         vector_store = get_qdrant_vector_store(user, collection_name)
#         logger.debug(f"Retrieving documents for vector_store_id: {vector_store_id}, collection: {collection_name}")
#         filter = Filter(
#             must=[
#                 FieldCondition(
#                     key="metadata.vector_store_id",
#                     match=MatchValue(value=vector_store_id)
#                 )
#             ]
#         )
#         search_result = vector_store.client.scroll(
#             collection_name=collection_name,
#             scroll_filter=filter,
#             limit=100,
#             with_vectors=False,
#             with_payload=True
#         )
#         documents = [
#             {"content": record.payload.get("page_content", ""), "metadata": record.payload.get("metadata", {})}
#             for record in search_result[0]
#         ]
#         logger.info(f"Retrieved {len(documents)} documents for vector_store_id: {vector_store_id}")
#         return documents
#     except Exception as e:
#         logger.error(f"Error retrieving documents for vector_store_id {vector_store_id}: {e}")
#         return []
def retrieve_documents_by_vector_id(vector_store_id: str, user, collection_name: str) -> List[Dict]:
    try:
        vector_store = get_qdrant_vector_store(user, collection_name)
        logger.debug(f"Retrieving documents for vector_store_id: {vector_store_id}, collection: {collection_name}")

        # Fetch documents from Qdrant where metadata.vector_store_id matches
        filter = Filter(
            must=[
                FieldCondition(
                    key="metadata.vector_store_id",
                    match=MatchValue(value=vector_store_id)
                )
            ]
        )
        search_result = vector_store.client.scroll(
            collection_name=collection_name,
            scroll_filter=filter,
            limit=100,
            with_vectors=False,
            with_payload=True
        )
        documents = [
            {"content": record.payload.get("page_content", ""), "metadata": record.payload.get("metadata", {})}
            for record in search_result[0]
        ]
        logger.info(f"Retrieved {len(documents)} documents from Qdrant for vector_store_id: {vector_store_id}")

        # Fetch documents from DocumentAccess
        vector_store_obj = get_object_or_404(VectorStore, id=vector_store_id, tenant=user.tenant)
        document_access = DocumentAccess.objects.filter(
            vector_store=vector_store_obj,
            document__tenant=user.tenant
        ).select_related('document')
        access_documents = []
        seen_document_ids = {doc["metadata"].get("document_id") for doc in documents}

        for access in document_access:
            doc = access.document
            if str(doc.id) not in seen_document_ids:
                # Fetch Qdrant data using document_id
                qdrant_docs = vector_store.client.scroll(
                    collection_name=collection_name,
                    scroll_filter=Filter(
                        must=[
                            FieldCondition(
                                key="metadata.document_id",
                                match=MatchValue(value=str(doc.id))
                            )
                        ]
                    ),
                    limit=100,
                    with_vectors=False,
                    with_payload=True
                )[0]
                if not qdrant_docs:
                    logger.warning(f"No Qdrant data found for document_id {doc.id} in DocumentAccess")
                    continue
                for record in qdrant_docs:
                    access_documents.append({
                        "content": record.payload.get("page_content", ""),
                        "metadata": record.payload.get("metadata", {})
                    })
                seen_document_ids.add(str(doc.id))

        documents.extend(access_documents)
        logger.info(f"Total documents retrieved (Qdrant + DocumentAccess): {len(documents)} for vector_store_id: {vector_store_id}")
        return documents
    except Exception as e:
        logger.error(f"Error retrieving documents for vector_store_id {vector_store_id}: {e}")
        return []


def insert_document_to_vectorstore(
    text: str,
    source_type: str,
    file_ext: str,
    document_id: str,
    user,
    collection_name: str,
    vector_store_id: str
):
    try:
        vector_store = get_qdrant_vector_store(user, collection_name)
        docs = smart_split_text(text, file_ext)
        logger.info(f"Splitting text into {len(docs)} chunks for document_id {document_id}")
        metadata = extract_metadata(text, file_ext)
        metadata.update({
            "document_id": document_id,
            "vector_store_id": vector_store_id,
            "source_type": source_type,
            "file_ext": file_ext,
            "tenant_id": str(user.tenant.id)
        })
        for doc in docs:
            if not hasattr(doc, "metadata") or not isinstance(doc.metadata, dict):
                doc.metadata = {}
            doc.metadata.update(metadata)
        vector_store.add_documents(docs, batch_size=64)
        logger.info(f"Inserted {len(docs)} chunks for document_id {document_id} into vector store {vector_store_id} in collection {collection_name}")
    except Exception as e:
        logger.error(f"Error inserting document {document_id} to vector store {vector_store_id}: {e}")
        raise

def delete_documents_by_vector_id(document_id: str, user, collection_name: str) -> bool:
    try:
        vector_store = get_qdrant_vector_store(user, collection_name)
        filter = Filter(must=[FieldCondition(key="metadata.document_id", match=MatchValue(value=document_id))])
        vector_store.client.delete(collection_name=collection_name, points_selector=filter)
        logger.info(f"Deleted documents with document_id {document_id} from collection {collection_name}")
        return True
    except Exception as e:
        logger.error(f"Error deleting documents from Qdrant: {e}")
        return False

def smart_split_text(text: str, file_ext: str) -> List[LangChainDocument]:
    try:
        if file_ext == ".json" or (text.startswith('{') and text.endswith('}')):
            return split_json_rows(text)
        elif file_ext in [".csv", ".xls", ".xlsx"]:
            return split_csv_rows(text)
        else:
            return split_unstructured_text(text)
    except Exception as e:
        logger.error(f"Error splitting text: {e}")
        return []

def split_json_rows(json_text: str) -> List[LangChainDocument]:
    try:
        data = json.loads(json_text)
        documents = []
        if isinstance(data, list):
            for i, item in enumerate(data):
                documents.append(LangChainDocument(page_content=str(item), metadata={"chunk": i + 1}))
        elif isinstance(data, dict):
            for key, value in data.items():
                documents.append(LangChainDocument(page_content=str(value), metadata={"source": key, "chunk": 1}))
        return documents
    except Exception as e:
        logger.error(f"Error splitting JSON rows: {e}")
        return []

def split_csv_rows(text: str, max_lines_per_chunk: int = 50) -> List[LangChainDocument]:
    try:
        lines = text.strip().split("\n")
        if not lines:
            return []
        header = lines[0]
        data_lines = lines[1:]
        chunks = []
        for i in range(0, len(data_lines), max_lines_per_chunk):
            chunk_text = "\n".join([header] + data_lines[i:i + max_lines_per_chunk])
            chunks.append(LangChainDocument(page_content=chunk_text, metadata={"chunk": i // max_lines_per_chunk + 1}))
        return chunks
    except Exception as e:
        logger.error(f"Error splitting CSV rows: {e}")
        return []

def split_unstructured_text(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[LangChainDocument]:
    try:
        if not text.strip():
            return []
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", " ", ""]
        )
        docs = splitter.create_documents([text])
        for i, doc in enumerate(docs):
            doc.metadata["chunk"] = i + 1
        return docs
    except Exception as e:
        logger.error(f"Error splitting unstructured text: {e}")
        return []

def extract_metadata(text: str, file_ext: str) -> dict:
    try:
        metadata = {}
        if file_ext in [".csv", ".xls", ".xlsx"]:
            data = json.loads(text) if file_ext == ".json" else text
            if isinstance(data, (dict, list)):
                if isinstance(data, list) and data and isinstance(data[0], dict):
                    metadata["columns"] = list(data[0].keys())
                elif isinstance(data, dict):
                    metadata["keys"] = list(data.keys())
        elif file_ext in [".pdf", ".docx", ".pptx", ".txt"]:
            doc = nlp(text[:3000])
            entities = {"ORG": [], "GPE": [], "DATE": []}
            for ent in doc.ents:
                if ent.label_ in entities:
                    entities[ent.label_].append(ent.text)
            for key, value in entities.items():
                if value:
                    metadata[key.lower()] = list(set(value))
        return metadata
    except Exception as e:
        logger.error(f"Error extracting metadata: {e}")
        return {}

def enrich_document(document_obj, file_text: str, file_ext: str):
    try:
        if not file_text.strip():
            return
        summary = summarize_context(file_text[:3000], document_obj.tenant.user)
        metadata = extract_metadata(file_text, file_ext)
        document_obj.summary = summary
        document_obj.keywords = metadata
        document_obj.save()
        logger.info(f"Enriched document {document_obj.id} with summary and metadata")
    except Exception as e:
        logger.error(f"Error enriching document {document_obj.id}: {e}")

def detect_alerts(document_obj, file_text: str):
    try:
        alert_keywords = ["contract expiry", "payment due", "breach of contract", "submission deadline", "invoice", "cancellation policy"]
        file_text_lower = file_text.lower()
        for keyword in alert_keywords:
            if keyword in file_text_lower:
                idx = file_text_lower.find(keyword)
                snippet = file_text[max(0, idx-100): idx+100]
                DocumentAlert.objects.create(
                    document=document_obj,
                    keyword=keyword,
                    snippet=snippet
                )
                logger.info(f"Detected alert for keyword '{keyword}' in document {document_obj.id}")
    except Exception as e:
        logger.error(f"Error detecting alerts for document {document_obj.id}: {e}")

def summarize_context(context_text: str, user=None) -> str:
    try:
        api_key = get_openai_api_key(user.id if user else 0)
        llm = ChatOpenAI(api_key=api_key, temperature=0)
        prompt = f"Summarize the following content concisely and accurately:\n\n{context_text}\n\nSummary:"
        response = llm.invoke(prompt)
        return response.content if hasattr(response, "content") else str(response)
    except Exception as e:
        logger.error(f"Error summarizing context: {e}")
        return ""

def summarize_batch(batch_docs: List[LangChainDocument]) -> str:
    try:
        combined_context = "\n".join([doc.page_content for doc in batch_docs])
        return summarize_context(combined_context)
    except Exception as e:
        logger.error(f"Error summarizing batch: {e}")
        return ""

def batch_documents(documents: List[LangChainDocument], batch_size: int = 5):
    for i in range(0, len(documents), batch_size):
        yield documents[i:i + batch_size]